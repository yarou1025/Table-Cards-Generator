import collections
import collections.abc
import pptx
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from copy import deepcopy
import pandas as pd
import copy

# Load the Excel file
excel_file = 'example.xlsx'
df = pd.read_excel(excel_file, index_col=False, header=None)

# Get the first two columns
df_name_title = df.iloc[:, :2]

# Load the PowerPoint file
pptx_file = 'template.pptx'
prs = Presentation(pptx_file)

# Get the first slide in the presentation
slide = prs.slides[0]

# Get the first shape on the slide
titles = [slide.shapes[2], slide.shapes[4]]
names = [slide.shapes[3], slide.shapes[5]]

# Get the source slide and its shapes
source_slide = prs.slides[0]
source_shapes = source_slide.shapes

# Loop through the shapes on the source slide
for num in range(len(df_name_title.index)):
    # Add a new slide to the presentation
    new_slide = prs.slides.add_slide(
        prs.slide_layouts[6])

    for shp in source_shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left = shp.left
            top = shp.top
            width = shp.width
            height = shp.height
            new_shape = new_slide.shapes.add_picture(
                "nycu.png", left, top, width, height)
            new_shape.rotation = shp.rotation

        else:
            if shp.has_text_frame:
                text_frame = shp.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        cur_text = run.text
                        if shp in titles:
                            new_text = df_name_title.iloc[num][1]
                        else:
                            new_text = df_name_title.iloc[num][0]
                        run.text = new_text
                        print(run.text, num)

            el = shp.element
            newel = copy.deepcopy(el)

            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')


# Save the modified presentation
prs.save('modified.pptx')
